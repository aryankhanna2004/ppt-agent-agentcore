#!/usr/bin/env python3
"""Dump the visual identity of a .pptx as JSON.

Usage: inspect_template.py <path/to/template.pptx>

The agent calls this before authoring a new deck so it can reuse the
template's slide layouts, theme colours, fonts, and placeholder geometry
without having to parse OOXML itself.

Output: single JSON object on stdout, e.g.

  {
    "path": "/opt/templates/default-template.pptx",
    "slide_size": {"width_in": 13.333, "height_in": 7.5},
    "layouts": [
      {"index": 0, "name": "Title Slide",
       "placeholders": [{"idx": 0, "type": "TITLE", "name": "Title 1",
                         "left_in": 0.8, "top_in": 2.6,
                         "width_in": 11.5, "height_in": 1.6}]}
    ],
    "theme_colors": {"accent1": "E87A22", "bg1": "0F1E3C", ...},
    "fonts": {"major": "Calibri", "minor": "Calibri"},
    "slide_count": 4
  }
"""
from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _in(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(Emu(emu).inches, 3)


def _placeholder_info(ph) -> dict[str, Any]:
    ph_format = ph.placeholder_format
    ptype = ph_format.type
    type_name = None
    if ptype is not None:
        type_name = str(ptype).split(".")[-1].split(" ")[0]
    return {
        "idx": ph_format.idx,
        "type": type_name,
        "name": ph.name,
        "left_in": _in(ph.left),
        "top_in": _in(ph.top),
        "width_in": _in(ph.width),
        "height_in": _in(ph.height),
    }


def _theme_info(prs: Presentation) -> tuple[dict[str, str], dict[str, str]]:
    colors: dict[str, str] = {}
    fonts: dict[str, str] = {}

    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
    except (IndexError, KeyError):
        return colors, fonts

    theme_xml = theme_part.blob
    from lxml import etree

    theme = etree.fromstring(theme_xml)
    scheme = theme.find(f"{{{NS_A}}}themeElements/{{{NS_A}}}clrScheme")
    if scheme is not None:
        for child in scheme:
            name = child.tag.split("}")[-1]
            srgb = child.find(f"{{{NS_A}}}srgbClr")
            sys_clr = child.find(f"{{{NS_A}}}sysClr")
            if srgb is not None:
                colors[name] = srgb.get("val")
            elif sys_clr is not None:
                colors[name] = sys_clr.get("lastClr") or sys_clr.get("val")

    font_scheme = theme.find(f"{{{NS_A}}}themeElements/{{{NS_A}}}fontScheme")
    if font_scheme is not None:
        for role in ("majorFont", "minorFont"):
            f = font_scheme.find(f"{{{NS_A}}}{role}/{{{NS_A}}}latin")
            if f is not None:
                fonts[role.replace("Font", "")] = f.get("typeface")

    return colors, fonts


def inspect(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    colors, fonts = _theme_info(prs)

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layouts.append(
            {
                "index": i,
                "name": layout.name,
                "placeholders": [_placeholder_info(ph) for ph in layout.placeholders],
            }
        )

    return {
        "path": str(path),
        "slide_size": {
            "width_in": _in(prs.slide_width),
            "height_in": _in(prs.slide_height),
        },
        "layouts": layouts,
        "theme_colors": colors,
        "fonts": fonts,
        "slide_count": len(prs.slides),
    }


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("usage: inspect_template.py <template.pptx>", file=sys.stderr)
        sys.exit(2)

    target = Path(sys.argv[1])
    if not target.exists():
        print(f"not found: {target}", file=sys.stderr)
        sys.exit(1)

    data = inspect(target)
    print(json.dumps(data, indent=2))
