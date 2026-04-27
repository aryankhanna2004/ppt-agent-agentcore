#!/usr/bin/env python3
"""Build a neutral default .pptx template from scratch.

This runs at Docker build time. The output `default-template.pptx` is baked
into `/opt/templates/` inside the harness image and is used whenever the
agent is asked to create a deck without a user-supplied template.

The template is intentionally generic (no third-party trademarks, logos, or
brand assets) so it is safe to redistribute under MIT. Organisations that
want their own style should drop their branded `.pptx` into
`/mnt/data/templates/` at runtime — see `inspect_template.py` and the
harness system prompt for how the agent picks it up.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0F, 0x1E, 0x3C)
ACCENT = RGBColor(0xE8, 0x7A, 0x22)
BODY = RGBColor(0x2A, 0x2A, 0x2A)
MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_text(tf, text, *, size, color, bold=False, name="Calibri"):
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_footer_bar(slide, color):
    from pptx.shapes.autoshape import Shape
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(7.1),
        Inches(13.333),
        Inches(0.4),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    title_slide = prs.slides.add_slide(blank)
    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY

    title_box = title_slide.shapes.add_textbox(
        Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.6)
    )
    _set_text(
        title_box.text_frame,
        "Presentation Title",
        size=54,
        color=WHITE,
        bold=True,
    )

    subtitle_box = title_slide.shapes.add_textbox(
        Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.8)
    )
    _set_text(
        subtitle_box.text_frame,
        "Subtitle / presenter / date",
        size=24,
        color=ACCENT,
    )

    _add_footer_bar(title_slide, ACCENT)

    content_slide = prs.slides.add_slide(blank)
    bg2 = content_slide.background.fill
    bg2.solid()
    bg2.fore_color.rgb = WHITE

    header_rect = content_slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9)
    )
    _set_text(
        header_rect.text_frame,
        "Section heading",
        size=32,
        color=NAVY,
        bold=True,
    )

    body_box = content_slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(12.0), Inches(5.2)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.text = "Bullet one — replace with real content."
    for extra in (
        "Bullet two — keep bullets short and parallel.",
        "Bullet three — one idea per line.",
    ):
        p = tf.add_paragraph()
        p.text = extra
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(24)
            run.font.color.rgb = BODY

    _add_footer_bar(content_slide, NAVY)

    section_slide = prs.slides.add_slide(blank)
    bg3 = section_slide.background.fill
    bg3.solid()
    bg3.fore_color.rgb = NAVY
    sec_box = section_slide.shapes.add_textbox(
        Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.4)
    )
    _set_text(
        sec_box.text_frame,
        "Section Divider",
        size=44,
        color=WHITE,
        bold=True,
    )
    _add_footer_bar(section_slide, ACCENT)

    closing_slide = prs.slides.add_slide(blank)
    bg4 = closing_slide.background.fill
    bg4.solid()
    bg4.fore_color.rgb = WHITE
    close_box = closing_slide.shapes.add_textbox(
        Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.4)
    )
    _set_text(
        close_box.text_frame,
        "Thank you",
        size=54,
        color=NAVY,
        bold=True,
    )
    qbox = closing_slide.shapes.add_textbox(
        Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.8)
    )
    _set_text(qbox.text_frame, "Questions?", size=28, color=ACCENT)
    _add_footer_bar(closing_slide, NAVY)

    prs.save(out_path)


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("default-template.pptx")
    build(out)
    print(f"Wrote {out} ({out.stat().st_size} bytes)")
